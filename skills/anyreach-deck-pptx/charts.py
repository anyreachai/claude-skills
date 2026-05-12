"""
Anyreach Deck (PPTX) — Native PPTX Chart Library.

PPTX-equivalent of the inline-SVG chart library in the PDF version. We
use python-pptx's native chart functionality (XL_CHART_TYPE) so the
recipient can edit the underlying data and series in PowerPoint.

CHARTS PROVIDED
  - donut(slide, ..., data)            — pie/donut for share breakdowns
  - stacked_area(slide, ..., data)     — multi-series area chart over time
  - bridge_bars(slide, ..., rows)      — horizontal bars for bridges (drawn
                                          as shapes, not native charts)
  - comparison_bars(slide, ..., rows)  — head-to-head bar comparison
                                          (drawn as shapes for fine-grained
                                          color-coding control)

WHY SOME CHARTS ARE SHAPES, NOT NATIVE CHARTS
The bridge_bars and comparison_bars are visually closer to "labeled
horizontal bars" than to a quantitative chart. Drawing them as shape
groups gives us full control over per-row coloring and label placement
(including the inline note text inside each bar). They remain fully
editable as shapes inside PowerPoint.

USAGE
    from lib.charts import donut, stacked_area, bridge_bars, comparison_bars
"""

from pptx.util import Pt, Emu
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import (
    XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION,
)
from pptx.dml.color import RGBColor

from .tokens import (
    C, FONT_DISPLAY, FONT_BODY, FONT_MONO, TYPE, SPACE, px_to_emu, rgb,
    LABEL_SPACING_HUNDREDTHS,
)
from .components import add_rect, add_textbox, rich_text, label, mono


# ============================================================
# DONUT (native pie chart)
# ============================================================

def donut(slide, x, y, size, data, hole_pct=62, center_label=None,
          center_value=None, on_cream=False):
    """Native PPTX donut chart with optional center label + center value.

    Args:
        x, y: top-left in px
        size: square dimension in px
        data: list of {'name': str, 'value': float, 'color': hex_or_token}
        hole_pct: inner hole as % of outer (default 62 — donut, not pie)
        center_label: small uppercase label inside the hole
        center_value: large serif value inside the hole

    The chart data is an editable PowerPoint chart. Recipients can change
    the values and the segments will update.
    """
    from pptx.chart.data import CategoryChartData

    chart_data = CategoryChartData()
    chart_data.categories = [d['name'] for d in data]
    chart_data.add_series('Series 1', [d['value'] for d in data])

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        px_to_emu(x), px_to_emu(y),
        px_to_emu(size), px_to_emu(size),
        chart_data,
    )
    chart = chart_shape.chart

    # Strip default chart chrome
    chart.has_title = False
    chart.has_legend = False

    plot = chart.plots[0]
    plot.has_data_labels = False

    # Color each segment
    series = plot.series[0]
    for i, d in enumerate(data):
        pt = series.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = rgb(d['color'])
        pt.format.line.fill.background()  # no border between segments

    # Center label / center value (overlay textboxes in the donut hole)
    # Note: native PPTX doughnut charts have a small internal margin
    # (typically ~5%) so the visual center of the donut is the same as the
    # geometric center of the bounding box. We use narrow, precisely
    # positioned textboxes so the text sits cleanly in the hole.
    if center_label or center_value:
        text_color_tok = 'cream' if not on_cream else 'ink'
        # Hole label width: ~48% of donut size (well inside the inner ring)
        text_w = int(size * 0.48)
        text_x = x + (size - text_w) / 2
        center_y = y + size / 2

        if center_label:
            rich_text(
                slide,
                text_x, center_y - 28,
                text_w, 22,
                runs=[center_label],
                font_size=TYPE['micro'],
                font_face=FONT_BODY,
                font_color='mutedOnInk' if not on_cream else 'mutedOnCream',
                align='center',
                bold=True,
                letter_spacing=LABEL_SPACING_HUNDREDTHS,
                line_spacing=1.0,
            )
        if center_value:
            rich_text(
                slide,
                text_x, center_y - 2,
                text_w, 60,
                runs=[center_value] if isinstance(center_value, str) else center_value,
                font_size=TYPE['metric_md'],
                font_face=FONT_DISPLAY,
                font_color=text_color_tok,
                align='center',
                line_spacing=1.0,
            )


def donut_legend(slide, x, y, w, data, on_cream=False, row_h=34):
    """Right-side legend for a donut: color swatch + name + value + pct.

    Designed to live next to a donut chart. Computes total internally for
    pct display if rows already include 'pct'; otherwise uses raw values.
    """
    line_tok = 'creamLine' if on_cream else 'inkLine'
    text_tok = 'ink' if on_cream else 'cream'
    muted_tok = 'mutedOnCream' if on_cream else 'mutedOnInk'

    total = sum(d['value'] for d in data)
    cy = y
    for d in data:
        # Color swatch
        add_rect(slide, x, cy + 8, 14, 14, fill_token=d['color'])
        # Name
        rich_text(
            slide, x + 26, cy + 4, w - 200, 24,
            runs=[d['name']],
            font_size=TYPE['body_sm'],
            font_face=FONT_BODY,
            font_color=text_tok,
        )
        # Value (mono, right-aligned-ish)
        val_str = d.get('value_label') or f'${d["value"]:.1f}M'
        rich_text(
            slide, x + w - 140, cy + 4, 80, 24,
            runs=[mono(val_str)],
            font_size=TYPE['body_sm'],
            font_face=FONT_MONO,
            font_color=text_tok,
            align='right',
        )
        # Percent
        pct = d.get('pct', round(d['value'] / total * 100))
        rich_text(
            slide, x + w - 56, cy + 4, 56, 24,
            runs=[mono(f'{pct}%')],
            font_size=TYPE['body_sm'],
            font_face=FONT_MONO,
            font_color=muted_tok,
            align='right',
        )
        # Hairline
        add_rect(slide, x, cy + row_h, w, 1, fill_token=line_tok)
        cy += row_h


# ============================================================
# STACKED AREA (native chart)
# ============================================================

def stacked_area(slide, x, y, w, h, data, series, on_cream=False):
    """Native multi-series stacked area chart.

    Args:
        data: list of dicts, one per x-axis tick. Each dict has a
              'label' (x category) and one key per series name.
              e.g. {'label': 'M1', 'BPO': 5, 'Delivery': 3, ...}
        series: list of {'key': str, 'name': str, 'color': hex_or_token}
                in the order they should stack (bottom → top).

    Recipients can edit the chart data in PowerPoint to change values.
    """
    chart_data = CategoryChartData()
    chart_data.categories = [d['label'] for d in data]
    for s in series:
        chart_data.add_series(
            s['name'],
            [d.get(s['key'], 0) for d in data],
        )

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.AREA_STACKED,
        px_to_emu(x), px_to_emu(y),
        px_to_emu(w), px_to_emu(h),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_title = False

    # Color series
    for i, s in enumerate(series):
        plot_series = chart.plots[0].series[i]
        plot_series.format.fill.solid()
        plot_series.format.fill.fore_color.rgb = rgb(s['color'])
        plot_series.format.line.fill.background()

    # Style axes
    text_color = rgb('mutedOnInk' if not on_cream else 'mutedOnCream')

    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(TYPE['caption'])
    cat_axis.tick_labels.font.color.rgb = text_color
    cat_axis.tick_labels.font.name = FONT_BODY

    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(TYPE['caption'])
    val_axis.tick_labels.font.color.rgb = text_color
    val_axis.tick_labels.font.name = FONT_MONO

    # Legend at the bottom
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(TYPE['caption'])
    chart.legend.font.color.rgb = text_color
    chart.legend.font.name = FONT_BODY


# ============================================================
# BRIDGE BARS (drawn as shapes for full control)
# ============================================================

def bridge_bars(slide, x, y, w, rows, max_value=None, on_cream=False,
                row_h=44, gap=10):
    """Horizontal bars showing a bridge / waterfall. Drawn as shape rects
    rather than a native chart, so we get the inline note text inside
    each bar and full per-row color control.

    Args:
        rows: list of dicts with keys:
            'label':  left-side row label
            'value':  numeric value (for bar length)
            'note':   optional inline text inside the bar
            'color':  bar fill (token or hex)
            'value_label': string to render on the right (e.g. "$3.0M")
        max_value: x-axis max (defaults to max of values * 1.1)
        row_h: bar height in px
        gap: vertical gap between bars

    Returns the y-coordinate just below the chart.
    """
    if max_value is None:
        max_value = max(r['value'] for r in rows) * 1.1

    text_tok = 'cream' if not on_cream else 'ink'
    muted_tok = 'mutedOnInk' if not on_cream else 'mutedOnCream'

    label_w = 180   # left column for row labels
    value_w = 100   # right column for value labels
    chart_w = w - label_w - value_w - 16  # bar drawing area

    cy = y
    for r in rows:
        # Label (left)
        rich_text(
            slide, x, cy + (row_h - 18) / 2, label_w, 22,
            runs=[r['label']],
            font_size=TYPE['body_sm'],
            font_face=FONT_BODY,
            font_color=muted_tok,
            line_spacing=1.4,
        )

        # Bar
        bar_w = chart_w * (r['value'] / max_value)
        bar_x = x + label_w
        add_rect(slide, bar_x, cy, bar_w, row_h,
                 fill_token=r.get('color', 'indigo'))

        # Inline note (inside the bar)
        if r.get('note'):
            rich_text(
                slide, bar_x + 12, cy + (row_h - 18) / 2, bar_w - 24, 22,
                runs=[r['note']],
                font_size=TYPE['caption'],
                font_face=FONT_BODY,
                font_color='cream',
                line_spacing=1.4,
            )

        # Value (right)
        rich_text(
            slide, x + label_w + chart_w + 16, cy + (row_h - 18) / 2,
            value_w, 22,
            runs=[mono(r.get('value_label', f'{r["value"]}'))],
            font_size=TYPE['body_sm'],
            font_face=FONT_MONO,
            font_color=text_tok,
            line_spacing=1.4,
        )

        cy += row_h + gap

    return cy


# ============================================================
# COMPARISON BARS (you vs benchmark, two-row groups)
# ============================================================

def comparison_bars(slide, x, y, w, benchmarks, on_cream=False,
                    row_h=20, group_gap=22, you_label='You',
                    them_label='Industry'):
    """Head-to-head bar pairs per metric. Each metric becomes a 2-row
    group: 'You' bar on top, 'Industry' bar below. Color-coded by who's
    "winning" — winner gets indigo, loser gets rose.

    Args:
        benchmarks: list of dicts with keys:
            'metric':  metric name
            'you':     your value (numeric)
            'them':    benchmark value (numeric)
            'higher_is_better': bool (default True)
            'unit':    optional unit string (e.g. '%')
        on_cream: bg context

    Returns: y-coordinate just below the chart.
    """
    text_tok = 'ink' if on_cream else 'cream'
    muted_tok = 'mutedOnCream' if on_cream else 'mutedOnInk'

    metric_w = 180
    chart_w = w - metric_w - 100
    value_w = 80

    cy = y
    for b in benchmarks:
        higher_better = b.get('higher_is_better', True)
        you_wins = (b['you'] > b['them']) if higher_better else (b['you'] < b['them'])

        max_val = max(b['you'], b['them']) * 1.0 or 1
        unit = b.get('unit', '')

        # Metric label (centered vertically across the 2-row group)
        rich_text(
            slide, x, cy + (row_h * 2 + 4 - 20) / 2, metric_w, 22,
            runs=[b['metric']],
            font_size=TYPE['body_sm'],
            font_face=FONT_BODY,
            font_color=text_tok,
            line_spacing=1.3,
            bold=True,
        )

        # Row 1: You
        you_color = 'indigo' if you_wins else 'rose'
        you_w = chart_w * (b['you'] / max_val)
        add_rect(slide, x + metric_w, cy, you_w, row_h,
                 fill_token=you_color)
        # 'You' tiny label inside or after bar
        rich_text(
            slide, x + metric_w + 6, cy + (row_h - 14) / 2, 80, 16,
            runs=[you_label],
            font_size=TYPE['micro'],
            font_face=FONT_BODY,
            font_color='cream',
            line_spacing=1.0,
            bold=True,
        )
        # Value
        rich_text(
            slide, x + metric_w + chart_w + 8, cy, value_w, row_h,
            runs=[mono(f'{b["you"]}{unit}')],
            font_size=TYPE['body_sm'],
            font_face=FONT_MONO,
            font_color=text_tok,
            line_spacing=1.4,
        )

        # Row 2: Them
        them_y = cy + row_h + 4
        them_color = 'rose' if you_wins else 'indigo'
        them_w = chart_w * (b['them'] / max_val)
        add_rect(slide, x + metric_w, them_y, them_w, row_h,
                 fill_token=them_color)
        rich_text(
            slide, x + metric_w + 6, them_y + (row_h - 14) / 2, 80, 16,
            runs=[them_label],
            font_size=TYPE['micro'],
            font_face=FONT_BODY,
            font_color='cream',
            line_spacing=1.0,
            bold=True,
        )
        rich_text(
            slide, x + metric_w + chart_w + 8, them_y, value_w, row_h,
            runs=[mono(f'{b["them"]}{unit}')],
            font_size=TYPE['body_sm'],
            font_face=FONT_MONO,
            font_color=muted_tok,
            line_spacing=1.4,
        )

        cy = them_y + row_h + group_gap

    return cy
