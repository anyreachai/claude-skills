"""
Anyreach Deck (PPTX) — Component Library.

Native-PPTX equivalent of the HTML component library in the PDF version.
Each function takes a slide and adds shapes/text to it at specific
coordinates. Coordinates are in "px" (matching the original 1400 x 900
design space) and converted to EMU internally.

NAMING CONVENTION
- blank_slide(...)     → creates a new slide with a background
- hero_block(...)      → places the title block on a slide
- section_header(...)  → places the standard section opener
- tile_grid(...)       → places a row/grid of tiles with hairline gaps
- metric_tile(...)     → returns a tile-content callable for tile_grid
- detail_tile(...)     → ditto, more elaborate variant
- qa_row(...)          → places one Q&A row
- scenario_card(...)   → places one scenario card
- callout(...)         → places a sensitivity / warning callout
- closing_tagline(...) → places the closing block

INLINE TEXT HELPERS — return Run dicts (text + style flags). Used inside
mixed-text helpers like rich_text():

    italic('hello')   -> {'text': 'hello', 'italic': True}
    mono('$15M')      -> {'text': '$15M', 'mono': True}

USAGE
    from lib.components import (
        blank_slide, hero_block, section_header,
        tile_grid, metric_tile, detail_tile,
        qa_row, scenario_card, callout, closing_tagline,
        italic, mono, label,
    )
"""

from copy import deepcopy

from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

from .tokens import (
    C, FONT_DISPLAY, FONT_BODY, FONT_MONO, TYPE, SPACE,
    PAGE_W, PAGE_H, px_to_emu, px_to_pt, rgb,
    LABEL_SPACING_HUNDREDTHS, LABEL_FONT_SIZE_PT, label_color,
)


# ============================================================
# INLINE TEXT HELPERS (return run-dicts, not HTML)
# ============================================================

def italic(text: str, color: str = None) -> dict:
    """Italic emphasis run. Use inside rich_text() lists.

        rich_text(slide, x, y, w, h, runs=[
            'Why ',
            italic('$15M', color='lime'),
            '.'
        ], font_size=TYPE['hero'], font_face=FONT_DISPLAY)
    """
    run = {'text': text, 'italic': True}
    if color:
        run['color'] = color
    return run


def mono(text: str, color: str = None) -> dict:
    """Monospace numeric run. Use inside rich_text() lists.

        rich_text(slide, x, y, w, h, runs=[
            'Total raise: ',
            mono('$15M'),
            '.'
        ], ...)
    """
    run = {'text': text, 'mono': True}
    if color:
        run['color'] = color
    return run


def label_run(text: str, color_token: str = None) -> dict:
    """Label-styled run (UPPERCASE, tracked, micro size). Less common as
    an inline run; usually use the standalone `label()` helper which
    creates a full text box.
    """
    run = {'text': text, 'is_label': True}
    if color_token:
        run['color'] = color_token
    return run


# ============================================================
# LOW-LEVEL DRAWING PRIMITIVES
# ============================================================

def add_rect(slide, x, y, w, h, fill_token, line_token=None):
    """Draw a filled rectangle. Coordinates and dimensions in px.
    fill_token / line_token: token names (e.g. 'cream') or 6-hex strings.
    line_token=None → no border.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        px_to_emu(x), px_to_emu(y),
        px_to_emu(w), px_to_emu(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_token)
    if line_token:
        shape.line.color.rgb = rgb(line_token)
        shape.line.width = Emu(0)  # hairline
    else:
        shape.line.fill.background()  # no outline
    shape.shadow.inherit = False  # disable PPT default shadow
    return shape


def add_textbox(slide, x, y, w, h):
    """Add a bare textbox at the given px coordinates. Returns the shape
    so the caller can populate text frame paragraphs/runs."""
    box = slide.shapes.add_textbox(
        px_to_emu(x), px_to_emu(y),
        px_to_emu(w), px_to_emu(h),
    )
    tf = box.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    return box


def _normalise_runs(runs):
    """Turn a mixed list of strings + run-dicts into a uniform list of
    run-dicts. Plain strings become {'text': str}."""
    normalised = []
    for r in runs:
        if isinstance(r, str):
            normalised.append({'text': r})
        else:
            normalised.append(dict(r))   # shallow copy so we can safely mutate
    return normalised


def _set_letter_spacing(run, hundredths: int):
    """Set the letter-spacing (tracking) of a run via raw XML.
    PowerPoint expresses spacing in 1/100 of a point on the rPr element.
    """
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(hundredths))


def rich_text(slide, x, y, w, h, runs, font_size, font_face=FONT_BODY,
              font_color='ink', bold=False, italic_default=False,
              line_spacing=1.15, align='left', vertical='top',
              letter_spacing=None):
    """Place mixed-style text in a textbox.

    Args:
        runs: list of strings and/or run-dicts (from italic(), mono()).
        font_size: default size (pt) — overridden per-run if 'size' set.
        font_face: default face — overridden per-run.
        font_color: default token name — overridden per-run.
        line_spacing: line height multiplier.
        letter_spacing: optional default tracking in 1/100 pt.

    Use for headlines, body paragraphs, metric values — any text that
    needs a single line/paragraph with multiple runs.
    """
    box = add_textbox(slide, x, y, w, h)
    tf = box.text_frame
    p = tf.paragraphs[0]
    if align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT
    p.line_spacing = line_spacing

    # Vertical anchor
    tf.vertical_anchor = {
        'top':    MSO_ANCHOR.TOP,
        'middle': MSO_ANCHOR.MIDDLE,
        'bottom': MSO_ANCHOR.BOTTOM,
    }.get(vertical, MSO_ANCHOR.TOP)

    for i, rdef in enumerate(_normalise_runs(runs)):
        r = p.add_run()
        r.text = rdef['text']
        f = r.font

        # Font face: mono > display/body > default
        if rdef.get('mono'):
            f.name = FONT_MONO
        elif rdef.get('face'):
            f.name = rdef['face']
        else:
            f.name = font_face

        # Italic
        if rdef.get('italic') or (italic_default and not rdef.get('italic') is False):
            f.italic = True
        elif italic_default:
            f.italic = True

        # Bold
        if rdef.get('bold') is not None:
            f.bold = rdef['bold']
        elif bold:
            f.bold = True

        # Size
        size_pt = rdef.get('size', font_size)
        f.size = Pt(size_pt)

        # Color
        color_tok = rdef.get('color', font_color)
        f.color.rgb = rgb(color_tok)

        # Tracking
        if rdef.get('letter_spacing') is not None:
            _set_letter_spacing(r, rdef['letter_spacing'])
        elif letter_spacing is not None:
            _set_letter_spacing(r, letter_spacing)

        # Label-style run (treated as inline label rather than full block)
        if rdef.get('is_label'):
            f.size = Pt(LABEL_FONT_SIZE_PT)
            f.bold = True
            r.text = rdef['text'].upper()
            _set_letter_spacing(r, LABEL_SPACING_HUNDREDTHS)

    return box


def label(slide, x, y, w, h, text, on_cream=True, color=None):
    """The omnipresent uppercase tracked-out label.

        label(slide, x=80, y=64, w=400, h=20,
              text='01 · USE OF FUNDS', on_cream=True)
    """
    color_tok = color or label_color(on_cream)
    box = add_textbox(slide, x, y, w, h)
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text.upper()
    f = r.font
    f.name = FONT_BODY
    f.size = Pt(LABEL_FONT_SIZE_PT)
    f.bold = True
    f.color.rgb = rgb(color_tok)
    _set_letter_spacing(r, LABEL_SPACING_HUNDREDTHS)
    return box


# ============================================================
# SLIDE WRAPPER
# ============================================================

def blank_slide(prs, background='cream'):
    """Add a new blank slide with a solid background color.

    Args:
        prs: pptx.Presentation
        background: token name ('cream' or 'ink' typically)

    Returns: the new Slide object.
    """
    layout = prs.slide_layouts[6]  # 'Blank' is index 6 in the default master
    slide = prs.slides.add_slide(layout)

    # Slide-level background fill via XML
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb(background)
    return slide


# ============================================================
# HERO BLOCK — page 1 only
# ============================================================

def hero_block(slide, eyebrow, headline_runs, subhead=None, byline=None,
               x=None, y=None, w=None,
               headline_size=None, on_cream=True):
    """The deck's opening title block.

    Args:
        eyebrow: small uppercase label
        headline_runs: a string or list of runs (use italic() for
                       emphasis, optionally with color='lime').
        subhead: optional supporting paragraph
        byline: optional small line of contact info
        x, y, w: position + width in px (default uses page padding)
        headline_size: pt size (default TYPE['hero'])
    """
    pad_x, pad_y = SPACE['page_pad_x'], SPACE['page_pad_y']
    x = x if x is not None else pad_x
    y = y if y is not None else pad_y
    w = w if w is not None else PAGE_W - 2 * pad_x
    h_size = headline_size or TYPE['hero']

    # Eyebrow
    label(slide, x, y, w, 26, eyebrow, on_cream=on_cream)

    # Headline (massive serif). Headline box height tracks font size so
    # we don't waste vertical space. h_size * 1.65 covers one full line
    # + descender buffer for a 135pt headline.
    if isinstance(headline_runs, str):
        headline_runs = [headline_runs]
    head_h = int(h_size * 1.65)
    rich_text(
        slide, x, y + 44, w, head_h,
        runs=headline_runs,
        font_size=h_size,
        font_face=FONT_DISPLAY,
        font_color='ink' if on_cream else 'cream',
        line_spacing=0.95,
    )

    cursor_y = y + 44 + head_h
    if subhead:
        # Subhead is body_lg=17pt at line_spacing 1.55, ~35px per line.
        # 3 lines + a 14px gap to byline ≈ 130 advance.
        rich_text(
            slide, x, cursor_y, min(w, 720), 130,
            runs=[subhead] if isinstance(subhead, str) else subhead,
            font_size=TYPE['body_lg'],
            font_face=FONT_BODY,
            font_color='mutedOnCream' if on_cream else 'mutedOnInk',
            line_spacing=1.55,
        )
        cursor_y += 130

    if byline:
        rich_text(
            slide, x, cursor_y, w, 28,
            runs=[byline],
            font_size=TYPE['body_xs'],
            font_face=FONT_BODY,
            font_color='mutedOnCream' if on_cream else 'mutedOnInk',
        )


# ============================================================
# SECTION HEADER — every non-hero page starts with this
# ============================================================

def section_header(slide, label_text, headline_runs, body=None,
                   on_cream=True, headline_size=None,
                   x=None, y=None, w=None):
    """Standard section opener: label + serif headline + optional body.

    Returns: the y-coordinate (px) just BELOW the header — convenient
    for placing the next content block.
    """
    pad_x, pad_y = SPACE['page_pad_x'], SPACE['page_pad_y']
    x = x if x is not None else pad_x
    y = y if y is not None else pad_y
    w = w if w is not None else PAGE_W - 2 * pad_x
    h_size = headline_size or TYPE['section_lg']

    # Label
    label(slide, x, y, w, 22, label_text, on_cream=on_cream)
    cursor_y = y + 32

    # Headline. Headline box must be tall enough for 2 lines of text at
    # the chosen size. line_spacing 1.02 → each line ≈ h_size * 1.36 px.
    # Two lines + a small safety buffer.
    if isinstance(headline_runs, str):
        headline_runs = [headline_runs]
    head_h = int(h_size * 1.5 * 2)
    rich_text(
        slide, x, cursor_y, w, head_h,
        runs=headline_runs,
        font_size=h_size,
        font_face=FONT_DISPLAY,
        font_color='ink' if on_cream else 'cream',
        line_spacing=1.02,
    )
    cursor_y += head_h

    if body:
        # Body box sized for 2 lines at body=13.5pt, line_spacing 1.65,
        # ≈ 30 px per line. 70px = 2 lines + buffer.
        body_h = 70
        rich_text(
            slide, x, cursor_y + 4, min(w, 820), body_h,
            runs=[body] if isinstance(body, str) else body,
            font_size=TYPE['body'],
            font_face=FONT_BODY,
            font_color='mutedOnCream' if on_cream else 'mutedOnInk',
            line_spacing=1.65,
        )
        cursor_y += body_h

    return cursor_y + 28


# ============================================================
# TILE GRID — workhorse layout
# ============================================================

def tile_grid(slide, x, y, w, h, tile_fns, columns, on_cream=True,
              gap_color=None):
    """Place a grid of tiles separated by hairline gaps. The trick is the
    same as the HTML version: a divider-colored backing rect, with
    individual tile rects placed inside leaving 1-2px gaps that reveal
    the divider color.

    Args:
        x, y, w, h: bounding box in px
        tile_fns: list of callables. Each takes (slide, tile_x, tile_y,
                  tile_w, tile_h) and draws the tile's content.
        columns: number of columns (rows derived from len(tile_fns))
        on_cream: bg context (controls divider + tile colors)
        gap_color: override divider color token

    Tip: build tile_fns with metric_tile() / detail_tile() / a custom
    lambda. Each function decides what to draw inside its rect.
    """
    n = len(tile_fns)
    rows = (n + columns - 1) // columns
    line_tok = gap_color or ('creamLine' if on_cream else 'inkLine')
    tile_bg_tok = 'cream' if on_cream else 'ink'

    # Backing rect (gap color)
    add_rect(slide, x, y, w, h, fill_token=line_tok)

    gap = SPACE['tile_gap']
    cell_w = (w - gap * (columns - 1)) / columns
    cell_h = (h - gap * (rows - 1)) / rows

    for i, fn in enumerate(tile_fns):
        r, c = divmod(i, columns)
        tx = x + c * (cell_w + gap)
        ty = y + r * (cell_h + gap)
        # Tile background rect
        add_rect(slide, tx, ty, cell_w, cell_h, fill_token=tile_bg_tok)
        # Content (caller draws into tx, ty, cell_w, cell_h with
        # SPACE['card_pad'] of padding implied)
        pad = SPACE['card_pad']
        fn(slide, tx + pad, ty + pad, cell_w - 2 * pad, cell_h - 2 * pad)


# ============================================================
# METRIC TILE — small label, big serif number, optional sub
# ============================================================

def metric_tile(label_text, value, sub=None, on_cream=True, size=None,
                value_color=None):
    """Returns a callable that draws a metric tile inside a tile_grid cell.

        metric_tile_row uses these. You can also pass them directly to
        tile_grid:

        tile_grid(slide, x, y, w, h, tile_fns=[
            metric_tile('Raise', '$15M', 'Series A · Q4 2026'),
            metric_tile('Runway', '24 mo', 'Base case'),
            ...
        ], columns=4, on_cream=True)
    """
    def draw(slide, x, y, w, h):
        text_color_tok = 'ink' if on_cream else 'cream'
        muted_tok = label_color(on_cream)
        v_size = size or TYPE['metric_lg']
        v_color = value_color or text_color_tok

        label(slide, x, y, w, 22, label_text, on_cream=on_cream)
        rich_text(
            slide, x, y + 36, w, int(v_size * 1.2),
            runs=[value if isinstance(value, str) else value],
            font_size=v_size,
            font_face=FONT_DISPLAY,
            font_color=v_color,
            line_spacing=1.0,
        )
        if sub:
            sub_y = y + 36 + int(v_size * 1.1) + 12
            rich_text(
                slide, x, sub_y, w, 60,
                runs=[sub] if isinstance(sub, str) else sub,
                font_size=TYPE['caption'],
                font_face=FONT_BODY,
                font_color=muted_tok,
                line_spacing=1.5,
            )
    return draw


def metric_tile_row(slide, x, y, w, h, metrics, on_cream=True):
    """Convenience: a row of metric dicts → tile_grid with one row.

    metrics: list of dicts {'label': str, 'value': str, 'sub': str}
    """
    tile_fns = [
        metric_tile(m['label'], m['value'], m.get('sub'),
                    on_cream=on_cream,
                    size=m.get('size'),
                    value_color=m.get('value_color'))
        for m in metrics
    ]
    tile_grid(slide, x, y, w, h, tile_fns, columns=len(metrics),
              on_cream=on_cream)


# ============================================================
# DETAIL TILE — title + big number + ratio + bullets
# ============================================================

def detail_tile(eyebrow, eyebrow_aux, title, value, value_caption,
                ratio, bullets, on_cream=True, accent='indigo'):
    """A more elaborate tile: eyebrow row + title + big metric + italic
    ratio + em-dash bullet list.

    Returns a callable for use inside tile_grid.
    """
    def draw(slide, x, y, w, h):
        text_color_tok = 'ink' if on_cream else 'cream'
        muted_tok = label_color(on_cream)
        line_tok = 'creamLine' if on_cream else 'inkLine'

        cursor = y

        # Row 1: eyebrow on left, mono aux on right
        label(slide, x, cursor, int(w * 0.5), 18, eyebrow, on_cream=on_cream)
        if eyebrow_aux:
            rich_text(
                slide, x + int(w * 0.5), cursor, int(w * 0.5), 18,
                runs=[mono(eyebrow_aux)],
                font_size=TYPE['caption'],
                font_color=muted_tok,
                align='right',
            )
        cursor += 28

        # Row 2: title (serif, mid)
        rich_text(
            slide, x, cursor, w, 40,
            runs=[title],
            font_size=TYPE['card_title'],
            font_face=FONT_DISPLAY,
            font_color=text_color_tok,
            line_spacing=1.15,
        )
        cursor += 48

        # Row 3: big metric value (own line, no inline caption)
        # 48pt font = ~64 px tall (96/72 * 48). Allow 72px to be safe.
        rich_text(
            slide, x, cursor, w, 76,
            runs=[value],
            font_size=TYPE['metric_lg'],
            font_face=FONT_DISPLAY,
            font_color=accent,
            line_spacing=1.0,
        )
        cursor += 72

        # Row 4: small caption beneath the value
        if value_caption:
            rich_text(
                slide, x, cursor, w, 22,
                runs=[value_caption],
                font_size=TYPE['caption'],
                font_face=FONT_BODY,
                font_color=muted_tok,
                line_spacing=1.3,
            )
            cursor += 24

        # Row 5: italic ratio
        rich_text(
            slide, x, cursor, w, 28,
            runs=[ratio],
            font_size=TYPE['caption'],
            font_face=FONT_BODY,
            font_color=muted_tok,
            italic_default=True,
            line_spacing=1.4,
        )
        cursor += 28

        # Hairline divider
        div_y = cursor + 10
        add_rect(slide, x, div_y, w, 1, fill_token=line_tok)
        cursor = div_y + 12

        # Bullets (em-dash + text)
        for b in bullets:
            rich_text(
                slide, x, cursor, w, 22,
                runs=[
                    {'text': '— ', 'color': accent},
                    {'text': b, 'color': muted_tok},
                ],
                font_size=TYPE['caption'],
                font_face=FONT_BODY,
                line_spacing=1.5,
            )
            cursor += 22

    return draw


# ============================================================
# Q&A ROW
# ============================================================

def qa_row(slide, x, y, w, number, question, answer_runs, on_cream=True,
           row_h=None):
    """Place one Q&A row. Returns the y-coordinate just below the row
    (including the divider) so the caller can chain rows vertically.

    Args:
        number: 1-indexed question number → renders as 'Q1', 'Q2'...
        question: bold question text
        answer_runs: string or list of runs
    """
    text_color_tok = 'ink' if on_cream else 'cream'
    muted_tok = label_color(on_cream)
    line_tok = 'creamLine' if on_cream else 'inkLine'

    # 'Q1' big serif
    rich_text(
        slide, x, y, 72, 40,
        runs=[f'Q{number}'],
        font_size=TYPE['card_title'],
        font_face=FONT_DISPLAY,
        font_color='indigo',
        line_spacing=1.0,
    )

    # Question + answer column
    qa_x = x + 96
    qa_w = w - 96

    rich_text(
        slide, qa_x, y, qa_w, 36,
        runs=[question],
        font_size=TYPE['body'],
        font_face=FONT_BODY,
        font_color=text_color_tok,
        bold=True,
        line_spacing=1.4,
    )

    # Estimate answer height by length
    if isinstance(answer_runs, str):
        answer_runs = [answer_runs]
    plain_len = sum(len(r) if isinstance(r, str) else len(r.get('text', ''))
                    for r in answer_runs)
    # body_sm at 12pt with line_spacing 1.65 → ~22 px per line.
    # Wrap point ≈ 110 chars/line at qa_w≈1144px (calibrated empirically).
    est_lines = max(1, plain_len // 110 + 1)
    ans_h = est_lines * 22

    rich_text(
        slide, qa_x, y + 36, qa_w, ans_h + 8,
        runs=answer_runs,
        font_size=TYPE['body_sm'],
        font_face=FONT_BODY,
        font_color=muted_tok,
        line_spacing=1.65,
    )

    total_h = max(80, 36 + ans_h + 22)
    if row_h is not None:
        total_h = row_h

    # Divider line at the bottom
    add_rect(slide, x, y + total_h, w, 1, fill_token=line_tok)

    return y + total_h + 22


# ============================================================
# SCENARIO CARD — bull/base/bear style
# ============================================================

def scenario_card(eyebrow, headline_value, headline_unit,
                  metric_a_label, metric_a_value,
                  metric_b_label, metric_b_value,
                  body, implication, accent, is_highlight=False):
    """Returns a callable for use inside tile_grid (intended for ink bg).

    The highlighted scenario gets a colored top border and a slightly
    raised background.
    """
    def draw(slide, x, y, w, h):
        # Override the default tile bg if this is the highlighted card.
        # The tile's own backing rect was already drawn by tile_grid; we
        # paint over it.
        bg_tok = 'inkSoft' if is_highlight else 'ink'
        # Subtract padding (tile_grid added card_pad on all sides)
        pad = SPACE['card_pad']
        outer_x = x - pad
        outer_y = y - pad
        outer_w = w + 2 * pad
        outer_h = h + 2 * pad
        add_rect(slide, outer_x, outer_y, outer_w, outer_h, fill_token=bg_tok)
        if is_highlight:
            add_rect(slide, outer_x, outer_y, outer_w, 3, fill_token='lime')

        cursor_y = y

        # Eyebrow + status dot
        label(slide, x, cursor_y, w - 16, 18, eyebrow, on_cream=False,
              color=accent)
        add_oval(slide, x + w - 8, cursor_y + 4, 8, 8, accent)

        cursor_y += 30

        # Headline value + unit
        rich_text(
            slide, x, cursor_y, w, 40,
            runs=[
                {'text': headline_value, 'face': FONT_DISPLAY,
                 'color': 'cream', 'size': 22.5},
                {'text': '  ' + headline_unit, 'color': 'mutedOnInk',
                 'size': 12},
            ],
            font_size=22.5,
            font_face=FONT_DISPLAY,
            line_spacing=1.1,
        )
        cursor_y += 54

        # Two-column metric row
        col_w = (w - 16) // 2
        for col_i, (lbl, val) in enumerate([
            (metric_a_label, metric_a_value),
            (metric_b_label, metric_b_value),
        ]):
            cx = x + col_i * (col_w + 16)
            label(slide, cx, cursor_y, col_w, 14, lbl, on_cream=False)
            rich_text(
                slide, cx, cursor_y + 16, col_w, 22,
                runs=[mono(val)],
                font_size=11.25,
                font_face=FONT_MONO,
                font_color='cream',
                line_spacing=1.2,
            )

        cursor_y += 56

        # Body paragraph
        rich_text(
            slide, x, cursor_y, w, 70,
            runs=[body] if isinstance(body, str) else body,
            font_size=TYPE['body_xs'],
            font_face=FONT_BODY,
            font_color='mutedOnInk',
            line_spacing=1.6,
        )
        cursor_y += 70

        # Hairline + italic implication line in accent
        add_rect(slide, x, cursor_y, w, 1, fill_token='inkLine')
        rich_text(
            slide, x, cursor_y + 12, w, 50,
            runs=[implication] if isinstance(implication, str) else implication,
            font_size=TYPE['body_xs'],
            font_face=FONT_BODY,
            font_color=accent,
            italic_default=True,
            line_spacing=1.6,
        )

    return draw


def add_oval(slide, x, y, w, h, fill_token):
    """Add a small filled oval (used for status dots in scenario cards)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        px_to_emu(x), px_to_emu(y),
        px_to_emu(w), px_to_emu(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_token)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


# ============================================================
# CALLOUT — sensitivity / warning / accent box
# ============================================================

def callout(slide, x, y, w, h, eyebrow, body_runs, accent='crimson',
            on_cream=True):
    """A visual callout — colored left border, soft fill, eyebrow,
    body text. Always crimson by default (warnings/sensitivities).

    Sizing tip: the height parameter is for the visible box. If passing
    a tall h, the body text will be top-aligned with empty space below.
    To get a tight callout, pass an h close to your estimated content
    height (rule of thumb: 60 + 18px per body line).
    """
    bg_tok = 'creamSoft' if on_cream else 'inkSoft'
    text_color_tok = 'ink' if on_cream else 'cream'

    # Soft fill background
    add_rect(slide, x, y, w, h, fill_token=bg_tok)
    # Colored left border (4px — wide enough to read clearly)
    add_rect(slide, x, y, 4, h, fill_token=accent)

    pad_x = 26
    pad_y = 22
    label(slide, x + pad_x, y + pad_y, w - 2 * pad_x, 18, eyebrow,
          on_cream=on_cream, color=accent)

    rich_text(
        slide, x + pad_x, y + pad_y + 30, w - 2 * pad_x, h - pad_y - 36,
        runs=[body_runs] if isinstance(body_runs, str) else body_runs,
        font_size=TYPE['body_sm'],
        font_face=FONT_BODY,
        font_color=text_color_tok,
        line_spacing=1.6,
    )


# ============================================================
# CLOSING TAGLINE
# ============================================================

def closing_tagline(slide, x, y, w, headline_runs, lines=None,
                    on_cream=False):
    """The closing block at the end of a deck. Centered serif tagline +
    small uppercase metadata lines below. Defaults to ink bg context.

    Args:
        headline_runs: string or list of runs (use italic + accent spans)
        lines: list of small metadata strings
    """
    line_tok = 'creamLine' if on_cream else 'inkLine'
    muted_tok = label_color(on_cream)
    text_color_tok = 'ink' if on_cream else 'cream'

    # Hairline above
    add_rect(slide, x, y, w, 1, fill_token=line_tok)

    # Centered tagline
    rich_text(
        slide, x, y + 30, w, 200,
        runs=[headline_runs] if isinstance(headline_runs, str) else headline_runs,
        font_size=36,
        font_face=FONT_DISPLAY,
        font_color=text_color_tok,
        align='center',
        line_spacing=1.15,
    )

    if lines:
        # First line: tracked uppercase
        ly = y + 30 + 110
        rich_text(
            slide, x, ly, w, 24,
            runs=[lines[0]],
            font_size=TYPE['caption'],
            font_face=FONT_BODY,
            font_color=muted_tok,
            align='center',
            letter_spacing=LABEL_SPACING_HUNDREDTHS,
            bold=True,
        )
        ly += 24
        for ln in lines[1:]:
            rich_text(
                slide, x, ly, w, 18,
                runs=[ln],
                font_size=TYPE['body_xs'],
                font_face=FONT_BODY,
                font_color=muted_tok,
                align='center',
            )
            ly += 18
